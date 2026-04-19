"""Generate the customer-facing Acme Corp QBR deck (8 slides).

Output: sample_inputs/slides.pptx
Voice: polished, slightly salesy — the deck WE'RE presenting. Deliberately contains
contradictions with notes.pdf to test the Gem's adversarial grounding.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).resolve().parent.parent / "sample_inputs" / "slides.pptx"

# Midnight Executive palette
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x22, 0x22, 0x2B)
MUTED = RGBColor(0x6B, 0x73, 0x85)
ACCENT = RGBColor(0xE8, 0xB5, 0x4A)  # warm gold for callouts

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background() if line is None else None
    if line is None:
        s.line.color.rgb = fill
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=CHARCOAL,
             font="Calibri", align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, color=CHARCOAL, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def page_number(slide, n, total=8):
    add_text(slide, Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3),
             f"{n} / {total}", size=9, color=MUTED, font="Calibri",
             align=PP_ALIGN.RIGHT)


def footer_brand(slide, dark=False):
    color = ICE if dark else MUTED
    add_text(slide, Inches(0.5), Inches(7.1), Inches(5), Inches(0.3),
             "CONFIDENTIAL  ·  ACME CORP QBR  ·  APRIL 2026",
             size=9, color=color, font="Calibri")


def left_accent_bar(slide, color=NAVY):
    # signature motif: thick navy bar down the left edge
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, color)


# ============================================================
def slide_1_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    # Gold accent block
    add_rect(s, Inches(0.8), Inches(2.6), Inches(0.5), Inches(0.12), ACCENT)

    add_text(s, Inches(0.8), Inches(2.9), Inches(11), Inches(1.2),
             "Acme Corp", size=64, bold=True, color=WHITE, font="Calibri")
    add_text(s, Inches(0.8), Inches(3.9), Inches(11), Inches(0.8),
             "Quarterly Business Review", size=32, color=ICE, font="Calibri")
    add_text(s, Inches(0.8), Inches(4.7), Inches(11), Inches(0.5),
             "Q1 2026 Review and Path Forward", size=20, italic=True,
             color=ICE, font="Calibri")

    add_text(s, Inches(0.8), Inches(6.6), Inches(11), Inches(0.3),
             "April 2026", size=12, color=ICE, font="Calibri")
    footer_brand(s, dark=True)
    page_number(s, 1)


def slide_2_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    left_accent_bar(s, NAVY)
    add_text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
             "Agenda", size=40, bold=True, color=NAVY, font="Calibri")
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
             "Six blocks, 60 minutes, discussion over slides.",
             size=14, italic=True, color=MUTED, font="Calibri")

    items = [
        ("01", "Value delivered this quarter"),
        ("02", "Usage and adoption trends"),
        ("03", "Support and reliability"),
        ("04", "Roadmap and 2026 vision"),
        ("05", "Renewal and partnership"),
        ("06", "Open discussion"),
    ]
    y = 2.3
    for num, label in items:
        # Numbered circle
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y),
                               Inches(0.55), Inches(0.55))
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.line.color.rgb = NAVY
        tf = c.text_frame
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.name = "Calibri"; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = WHITE
        add_text(s, Inches(1.7), Inches(y + 0.08), Inches(10), Inches(0.5),
                 label, size=20, color=CHARCOAL, font="Calibri")
        y += 0.7

    footer_brand(s)
    page_number(s, 2)


def slide_3_value(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    left_accent_bar(s, NAVY)
    add_text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
             "Value delivered this quarter", size=36, bold=True,
             color=NAVY, font="Calibri")
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
             "A strong Q1 on every dimension we measure.",
             size=14, italic=True, color=MUTED, font="Calibri")

    # Three big stat cards
    stats = [
        ("2.4x", "ROI demonstrated through\nreduced analytics team time"),
        ("+38%", "Quarter-over-quarter\nquery growth"),
        ("12", "New dashboards deployed\nto production"),
    ]
    card_w = Inches(3.7); card_h = Inches(3.2)
    gap = Inches(0.3)
    start_x = (SLIDE_W - (card_w * 3 + gap * 2)) / 2
    y = Inches(2.4)
    for i, (big, small) in enumerate(stats):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, y, card_w, card_h, ICE)
        # colored top band
        add_rect(s, x, y, card_w, Inches(0.15), NAVY)
        add_text(s, x + Inches(0.3), y + Inches(0.4), card_w - Inches(0.6),
                 Inches(1.6), big, size=72, bold=True, color=NAVY,
                 font="Calibri", align=PP_ALIGN.LEFT)
        add_text(s, x + Inches(0.3), y + Inches(2.0), card_w - Inches(0.6),
                 Inches(1.0), small, size=14, color=CHARCOAL,
                 font="Calibri", align=PP_ALIGN.LEFT)

    footer_brand(s)
    page_number(s, 3)


def slide_4_usage(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    left_accent_bar(s, NAVY)
    add_text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
             "Usage and adoption", size=36, bold=True, color=NAVY,
             font="Calibri")
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
             "Growth across every cohort we track.",
             size=14, italic=True, color=MUTED, font="Calibri")

    # Left: key numbers
    add_bullets(s, Inches(0.9), Inches(2.5), Inches(6), Inches(4.2), [
        "Monthly active users up 12% QoQ",
        "Power user cohort growing",
        "Query volume up 38% QoQ — broadening, not just deepening",
        "Dashboard creation velocity at an all-time high",
        "Self-serve adoption expanding beyond the core data team",
    ], size=18)

    # Right: ASCII-ish bar chart (native shapes)
    chart_x = Inches(7.3); chart_y = Inches(2.5)
    add_rect(s, chart_x, chart_y, Inches(5.2), Inches(4.2), ICE)
    add_text(s, chart_x + Inches(0.3), chart_y + Inches(0.2),
             Inches(4.6), Inches(0.4),
             "Monthly Active Users — trailing 4 quarters",
             size=12, bold=True, color=NAVY, font="Calibri")
    bar_heights = [1.8, 2.1, 2.5, 2.8]  # arbitrary growth bars
    bar_labels = ["Q2'25", "Q3'25", "Q4'25", "Q1'26"]
    base_y = chart_y + Inches(3.6)
    base_x = chart_x + Inches(0.5)
    bar_w = Inches(0.9); bar_gap = Inches(0.25)
    for i, h in enumerate(bar_heights):
        x = base_x + (bar_w + bar_gap) * i
        bh = Inches(h * 0.9)
        add_rect(s, x, base_y - bh, bar_w, bh, NAVY)
        add_text(s, x - Inches(0.05), base_y + Inches(0.05),
                 bar_w + Inches(0.1), Inches(0.3),
                 bar_labels[i], size=10, color=MUTED, font="Calibri",
                 align=PP_ALIGN.CENTER)

    footer_brand(s)
    page_number(s, 4)


def slide_5_support(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    left_accent_bar(s, NAVY)
    add_text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
             "Support and reliability", size=36, bold=True,
             color=NAVY, font="Calibri")
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
             "A platform Acme can build on.",
             size=14, italic=True, color=MUTED, font="Calibri")

    stats = [
        ("99.2%", "Uptime delivered"),
        ("18 hrs", "Average ticket\nresolution"),
        ("0", "Data incidents"),
    ]
    card_w = Inches(3.7); card_h = Inches(3.2)
    gap = Inches(0.3)
    start_x = (SLIDE_W - (card_w * 3 + gap * 2)) / 2
    y = Inches(2.4)
    for i, (big, small) in enumerate(stats):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, y, card_w, card_h, ICE)
        add_rect(s, x, y, card_w, Inches(0.15), NAVY)
        add_text(s, x + Inches(0.3), y + Inches(0.4), card_w - Inches(0.6),
                 Inches(1.8), big, size=64, bold=True, color=NAVY,
                 font="Calibri", align=PP_ALIGN.LEFT)
        add_text(s, x + Inches(0.3), y + Inches(2.1), card_w - Inches(0.6),
                 Inches(1.0), small, size=14, color=CHARCOAL,
                 font="Calibri", align=PP_ALIGN.LEFT)

    footer_brand(s)
    page_number(s, 5)


def slide_6_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    left_accent_bar(s, NAVY)
    add_text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
             "2026 Roadmap", size=36, bold=True, color=NAVY, font="Calibri")
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
             "Shared roadmap: building the next phase together.",
             size=16, italic=True, color=MUTED, font="Calibri")

    items = [
        ("H1", "Advanced anomaly detection",
         "Real-time statistical and ML-driven outlier surfacing across ingested streams."),
        ("H1", "Multi-region data residency",
         "EU + APAC primary regions with compliance-aligned routing."),
        ("H2", "Native LLM integration",
         "First-class primitives for semantic search, retrieval, and grounded Q&A."),
        ("H2", "Expanded EMEA coverage",
         "Regional support coverage and a named CSM for the EMEA business unit."),
    ]
    y = 2.4
    for tag, title, body in items:
        # half-year tag
        tag_box = add_rect(s, Inches(0.9), Inches(y), Inches(0.8),
                           Inches(0.45), NAVY)
        tf = tag_box.text_frame
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = tag
        r.font.name = "Calibri"; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = WHITE
        add_text(s, Inches(1.9), Inches(y - 0.02),
                 Inches(10.8), Inches(0.4),
                 title, size=18, bold=True, color=CHARCOAL, font="Calibri")
        add_text(s, Inches(1.9), Inches(y + 0.4),
                 Inches(10.8), Inches(0.6),
                 body, size=12, color=MUTED, font="Calibri")
        y += 1.1

    footer_brand(s)
    page_number(s, 6)


def slide_7_renewal(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    left_accent_bar(s, NAVY)
    add_text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
             "Renewal and partnership", size=36, bold=True,
             color=NAVY, font="Calibri")
    add_text(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
             "A proposal designed for the next phase of the partnership.",
             size=14, italic=True, color=MUTED, font="Calibri")

    items = [
        ("Proposed",
         "3-year renewal with locked-in rates",
         "Price certainty through 2029. No annual increases."),
        ("Expanded access",
         "EMEA business unit rollout",
         "Extend platform coverage to the EMEA BU, with regional support."),
        ("Executive sponsor program",
         "Named executive sponsors on both sides",
         "Quarterly executive cadence; priority access for roadmap input."),
    ]
    y = 2.3
    for label, title, body in items:
        add_rect(s, Inches(0.9), Inches(y), Inches(11.5), Inches(1.3), ICE)
        add_rect(s, Inches(0.9), Inches(y), Inches(0.15), Inches(1.3), NAVY)
        add_text(s, Inches(1.25), Inches(y + 0.15),
                 Inches(2.8), Inches(0.35),
                 label.upper(), size=10, bold=True, color=NAVY, font="Calibri")
        add_text(s, Inches(1.25), Inches(y + 0.45),
                 Inches(11), Inches(0.45),
                 title, size=18, bold=True, color=CHARCOAL, font="Calibri")
        add_text(s, Inches(1.25), Inches(y + 0.85),
                 Inches(11), Inches(0.4),
                 body, size=12, color=MUTED, font="Calibri")
        y += 1.45

    footer_brand(s)
    page_number(s, 7)


def slide_8_discussion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0.8), Inches(2.3), Inches(0.5), Inches(0.12), ACCENT)

    add_text(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.2),
             "Open discussion", size=56, bold=True, color=WHITE, font="Calibri")
    add_text(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(2.5),
             "What would make the next 12 months\na success for Acme?",
             size=32, italic=True, color=ICE, font="Calibri")

    footer_brand(s, dark=True)
    page_number(s, 8)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_title(prs)
    slide_2_agenda(prs)
    slide_3_value(prs)
    slide_4_usage(prs)
    slide_5_support(prs)
    slide_6_roadmap(prs)
    slide_7_renewal(prs)
    slide_8_discussion(prs)

    prs.save(OUT)
    print(f"Wrote: {OUT} ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
